"""Markdown -> PPTX exporter (pip-native, python-pptx).

edudoc-owned. References skills/skills-main/pptx-manipulation (python-pptx usage)
without copying it. Same content contract as the other exporters: it reads a
Markdown file (e.g. from ComposedReport.to_markdown()) and builds slides.

Mapping:
- first heading            -> title slide
- each following heading    -> a new content slide (title = heading text)
- paragraphs / list items   -> bullet lines; leading □ / ○ / ― / ※ set bullet level
- numeric tables            -> a native chart slide (clustered column)
- other tables              -> a table slide
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from .export_base import BaseExporter, ExportResult
from .markdown_blocks import Heading, ListBlock, Paragraph, Run, Table, parse_markdown

_KOREAN_FONT = "Malgun Gothic"
_MARKER_LEVEL = {"□": 0, "○": 1, "―": 2, "※": 3}
# derived total columns/rows are dropped from a chart so quarter bars stay comparable
_TOTAL_LABELS = {"합계", "총계", "소계", "계"}
_NUMBER_RE = re.compile(r"^-?\d+(?:\.\d+)?$")


class PptxExporter(BaseExporter):
    """Markdown file -> PPTX (python-pptx)."""

    supported_ext = (".pptx",)

    def __init__(self, korean_font: str = _KOREAN_FONT, include_charts: bool = False) -> None:
        self.korean_font = korean_font
        # charts are opt-in: a plain deck stays written (numeric tables render as tables)
        self.include_charts = include_charts

    def export(self, markdown_path: Path, output_path: Path) -> ExportResult:
        markdown_path = Path(markdown_path)
        output_path = Path(output_path)
        meta = {"exporter": self.name, "requires_optional_tool": False, "stabilized": True}

        if not self.can_export(output_path):
            return ExportResult(source=markdown_path, output=output_path, ok=False,
                                error=f"Unsupported output extension: {output_path.suffix} "
                                      f"(supported: {sorted(self.supported_ext)})", meta=meta,
                                error_code="export_unsupported_extension")
        if not markdown_path.exists():
            return ExportResult(source=markdown_path, output=output_path, ok=False,
                                error=f"Markdown source does not exist: {markdown_path}", meta=meta,
                                error_code="export_source_missing")

        try:
            blocks = parse_markdown(markdown_path.read_text(encoding="utf-8"))
            prs = Presentation()
            self._build(prs, blocks)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))
            return ExportResult(source=markdown_path, output=output_path, ok=True,
                                meta={**meta, "slides": len(prs.slides._sldIdLst)})
        except Exception as e:  # noqa: BLE001 - structured failure
            return ExportResult(source=markdown_path, output=output_path, ok=False,
                                error=repr(e), meta={**meta, "stabilized": False},
                                error_code="export_failed")

    # --- build ----------------------------------------------------------

    def _build(self, prs, blocks) -> None:  # noqa: ANN001
        title_done = False
        pending_title = None  # heading text awaiting its slide (so a table can own it)
        body_tf = None        # text frame of the current content slide

        def ensure_body():
            nonlocal body_tf, pending_title
            if body_tf is None:
                body_tf = self._content_slide(prs, pending_title or "")
                pending_title = None
            return body_tf

        for block in blocks:
            if isinstance(block, Heading):
                if not title_done:
                    self._title_slide(prs, _runs_text(block.runs))
                    title_done = True
                else:
                    pending_title = _runs_text(block.runs)
                body_tf = None
            elif isinstance(block, Paragraph):
                text = _runs_text(block.runs)
                if text:
                    self._add_bullet(ensure_body(), text)
            elif isinstance(block, ListBlock):
                tf = ensure_body()
                for item in block.items:
                    self._add_bullet(tf, _runs_text(item), base_level=1)
            elif isinstance(block, Table):
                title = pending_title or ""
                chart = _chartable_data(block) if self.include_charts else None
                if chart is not None:
                    self._chart_slide(prs, title, chart)
                else:
                    self._table_slide(prs, title, block)
                body_tf = None
                pending_title = None

        if len(prs.slides._sldIdLst) == 0:
            self._title_slide(prs, "(빈 문서)")

    def _title_slide(self, prs, title: str):  # noqa: ANN001
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        self._apply_font(slide.shapes.title.text_frame)
        return slide

    def _content_slide(self, prs, title: str):  # noqa: ANN001
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        self._apply_font(slide.shapes.title.text_frame)
        body = slide.placeholders[1].text_frame
        body.clear()
        return body

    def _add_bullet(self, body_tf, text: str, base_level: int = 0) -> None:  # noqa: ANN001
        level = base_level
        marker = text[:1]
        if marker in _MARKER_LEVEL:
            level = _MARKER_LEVEL[marker]
            text = text[1:].strip()
        # first paragraph of a fresh text frame is empty and reusable
        first = body_tf.paragraphs[0]
        para = first if (len(body_tf.paragraphs) == 1 and not first.runs) else body_tf.add_paragraph()
        para.text = text
        para.level = min(level, 4)
        for run in para.runs:
            run.font.name = self.korean_font
            run.font.size = Pt(18 if level == 0 else 16)

    def _table_slide(self, prs, title: str, table: Table) -> None:  # noqa: ANN001
        rows_data = ([table.header] if table.header else []) + table.rows
        if not rows_data:
            return
        ncols = max(len(r) for r in rows_data)
        nrows = len(rows_data)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        if title and slide.shapes.title is not None:
            slide.shapes.title.text = title
            self._apply_font(slide.shapes.title.text_frame)
        gt = slide.shapes.add_table(
            nrows, ncols, Inches(0.5), Inches(1.6), Inches(9), Inches(0.4 * nrows)
        ).table
        for r, row in enumerate(rows_data):
            for c in range(ncols):
                cell = gt.cell(r, c)
                cell.text = _runs_text(row[c]) if c < len(row) else ""
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = self.korean_font
                        run.font.size = Pt(12)

    def _chart_slide(self, prs, title: str, chart_data: dict) -> None:  # noqa: ANN001
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        if title and slide.shapes.title is not None:
            slide.shapes.title.text = title
            self._apply_font(slide.shapes.title.text_frame)

        cdata = CategoryChartData()
        cdata.categories = chart_data["categories"]
        for name, values in chart_data["series"]:
            cdata.add_series(name, values)

        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), Inches(1.6), Inches(9), Inches(5), cdata,
        )
        chart = frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.font.name = self.korean_font
        chart.font.size = Pt(12)

    def _apply_font(self, text_frame) -> None:  # noqa: ANN001
        for para in text_frame.paragraphs:
            for run in para.runs:
                run.font.name = self.korean_font


def _cell_text(cell: list[Run]) -> str:
    return _runs_text(cell)


def _to_number(text: str) -> float | None:
    cleaned = text.replace(",", "").replace(" ", "").strip()
    return float(cleaned) if _NUMBER_RE.match(cleaned) else None


def _is_total_label(text: str) -> bool:
    t = text.strip()
    return t in _TOTAL_LABELS or t.startswith("총")


def _chartable_data(table: Table) -> dict | None:
    """Return {categories, series:[(name, values)]} if the table is numeric.

    Shape assumed: header = [라벨열, 범주1, 범주2, ...]; each data row = [계열명, 값, 값, ...].
    Derived total columns/rows (합계·총계·… or names starting with '총') are dropped so the
    chart compares like-for-like. Returns None when the table is not chartable.
    """
    if not table.header or not table.rows:
        return None

    header = [_cell_text(c) for c in table.header]
    if len(header) < 2:
        return None

    # category columns = non-first header cells that are not totals
    cat_cols = [i for i in range(1, len(header)) if not _is_total_label(header[i])]
    if not cat_cols:
        return None

    categories = [header[i] for i in cat_cols]
    series: list[tuple[str, list[float]]] = []
    for row in table.rows:
        cells = [_cell_text(c) for c in row]
        if not cells:
            continue
        name = cells[0]
        if _is_total_label(name):
            continue  # skip total rows
        values = [_to_number(cells[i]) if i < len(cells) else None for i in cat_cols]
        if any(v is None for v in values):
            return None  # any non-numeric data cell -> not a chart, render as table
        series.append((name, values))

    if not series:
        return None
    return {"categories": categories, "series": series}


def _runs_text(runs: list[Run]) -> str:
    return "".join(r.text for r in runs).strip()
